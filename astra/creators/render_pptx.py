"""
Render a deck artifact to PowerPoint (.pptx) via python-pptx.

Why offer .pptx alongside .pdf:
- Investors and partners often want to edit / annotate the deck.
  PDFs are read-only; PPTX lets them rearrange, swap a number, or
  paste their own logo before forwarding.
- Some procurement / RFP responses require editable formats.
- If the founder wants to make minor tweaks without a full
  regeneration cycle, PPTX is the faster path.

The .pptx is produced from the same JSON content that drives the
PDF render — same kit, same brand colors, same slide schema. So
PDF and PPTX renders of the same artifact should look equivalent
(modulo PDF-only effects like CSS gradients).
"""

from __future__ import annotations

import io
import logging
import os
from typing import Any

import boto3
from botocore.client import Config

from astra.creators.kits import load_kit
from astra.creators.store import get_artifact, update_artifact_render_key

logger = logging.getLogger(__name__)


def _r2_client():
    return boto3.client(
        "s3",
        endpoint_url=os.environ["R2_ENDPOINT"],
        aws_access_key_id=os.environ["R2_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["R2_SECRET_ACCESS_KEY"],
        config=Config(signature_version="s3v4"),
        region_name="auto",
    )


def _r2_bucket() -> str:
    return os.environ.get("R2_BUCKET", "astra-backups")


def _hex_to_rgb(hexstr: str) -> tuple[int, int, int]:
    """Convert "#RRGGBB" or "RRGGBB" to (r, g, b) tuple."""
    h = (hexstr or "").lstrip("#")
    if len(h) != 6:
        return (0, 0, 0)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _build_pptx(content: dict[str, Any], kit) -> bytes:
    """Build a .pptx file from deck JSON content. Returns bytes."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        raise RuntimeError(
            "python-pptx not installed; add it to pyproject.toml deps"
        ) from e

    colors = kit.colors or {}
    primary = RGBColor(*_hex_to_rgb(colors.get("primary", "#0A1F3D")))
    secondary = RGBColor(*_hex_to_rgb(colors.get("secondary", "#FF6B35")))
    surface = RGBColor(*_hex_to_rgb(colors.get("surface", "#FFFFFF")))
    ink = RGBColor(*_hex_to_rgb(colors.get("ink", "#0F0F0F")))
    muted = RGBColor(*_hex_to_rgb(colors.get("muted", "#6B6B6B")))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    fonts = kit.fonts or {}
    display_font = ((fonts.get("display") or {}).get("family")) or "Helvetica"
    body_font = ((fonts.get("body") or {}).get("family")) or "Helvetica"

    # 16:9 — python-pptx widescreen default is 13.333" x 7.5"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    slides = content.get("slides", []) or []
    deck_title = content.get("title", "")
    deck_subtitle = content.get("subtitle", "")

    company_name = kit.name
    output_cfg = (kit.brand.get("output") or {}) or {}
    footer_enabled = output_cfg.get("slide_footer", True)
    slide_numbers = output_cfg.get("slide_numbers", True)

    def _add_textbox(slide, *, left, top, width, height, text,
                     font, size_pt, color, bold=False, align=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        # First paragraph — set text
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = text or ""
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color
        run.font.bold = bold
        return tb

    def _add_bullets(slide, *, left, top, width, height, bullets,
                     font, size_pt, color):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = f"•  {b}"
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.color.rgb = color

    def _fill(slide, color: RGBColor) -> None:
        # Background fill on the slide
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    for idx, s in enumerate(slides, 1):
        stype = s.get("type", "content")
        slide = prs.slides.add_slide(blank_layout)

        # Decide background
        if stype in ("cover", "section", "close"):
            _fill(slide, primary)
            text_color = white
        else:
            _fill(slide, surface)
            text_color = ink

        # Margins
        left_m = Inches(0.8)
        top_m = Inches(0.7)
        usable_w = Inches(11.7)
        usable_h = Inches(6.0)

        if stype == "cover":
            _add_textbox(
                slide,
                left=left_m, top=Inches(2.5),
                width=usable_w, height=Inches(2.0),
                text=s.get("title") or deck_title,
                font=display_font, size_pt=60, color=text_color, bold=True,
            )
            sub = s.get("subtitle") or deck_subtitle
            if sub:
                _add_textbox(
                    slide,
                    left=left_m, top=Inches(4.6),
                    width=usable_w, height=Inches(1.0),
                    text=sub, font=body_font, size_pt=24,
                    color=RGBColor(0xE0, 0xE0, 0xE0),
                )

        elif stype == "section":
            _add_textbox(
                slide,
                left=left_m, top=Inches(3.0),
                width=usable_w, height=Inches(1.5),
                text=s.get("title", ""),
                font=display_font, size_pt=52, color=text_color, bold=True,
            )
            if s.get("subtitle"):
                _add_textbox(
                    slide,
                    left=left_m, top=Inches(4.4),
                    width=usable_w, height=Inches(1.0),
                    text=s["subtitle"], font=body_font, size_pt=22,
                    color=RGBColor(0xE0, 0xE0, 0xE0),
                )

        elif stype == "data":
            if s.get("heading"):
                _add_textbox(
                    slide, left=left_m, top=top_m,
                    width=usable_w, height=Inches(0.8),
                    text=s["heading"], font=display_font,
                    size_pt=24, color=ink, bold=True,
                )
            # Hero number — title is the value
            _add_textbox(
                slide,
                left=left_m, top=Inches(2.0),
                width=usable_w, height=Inches(2.5),
                text=s.get("title", ""),
                font=display_font, size_pt=110, color=primary, bold=True,
            )
            if s.get("body_md"):
                _add_textbox(
                    slide,
                    left=left_m, top=Inches(4.8),
                    width=usable_w, height=Inches(1.6),
                    text=s["body_md"], font=body_font,
                    size_pt=18, color=ink,
                )

        elif stype == "quote":
            _add_textbox(
                slide,
                left=left_m, top=Inches(2.0),
                width=usable_w, height=Inches(3.5),
                text=f"\u201c{s.get('body_md') or s.get('heading','')}\u201d",
                font=display_font, size_pt=32, color=ink,
            )
            if s.get("subtitle"):
                _add_textbox(
                    slide,
                    left=left_m, top=Inches(5.3),
                    width=usable_w, height=Inches(0.6),
                    text=f"\u2014 {s['subtitle']}",
                    font=body_font, size_pt=16, color=muted,
                )

        elif stype == "close":
            _add_textbox(
                slide,
                left=left_m, top=Inches(2.4),
                width=usable_w, height=Inches(1.2),
                text=s.get("title") or "Let's talk.",
                font=display_font, size_pt=48, color=text_color, bold=True,
            )
            if s.get("body_md"):
                _add_textbox(
                    slide,
                    left=left_m, top=Inches(3.9),
                    width=usable_w, height=Inches(2.5),
                    text=s["body_md"],
                    font=body_font, size_pt=20,
                    color=RGBColor(0xE5, 0xE5, 0xE5),
                )
            if s.get("bullets"):
                _add_bullets(
                    slide, left=left_m, top=Inches(4.9),
                    width=usable_w, height=Inches(1.8),
                    bullets=s["bullets"],
                    font=body_font, size_pt=18,
                    color=RGBColor(0xE5, 0xE5, 0xE5),
                )

        else:
            # content slide
            cur_top = top_m
            if s.get("title"):
                _add_textbox(
                    slide, left=left_m, top=cur_top,
                    width=usable_w, height=Inches(1.0),
                    text=s["title"], font=display_font,
                    size_pt=32, color=ink, bold=True,
                )
                cur_top = Inches(1.7)
            if s.get("subtitle"):
                _add_textbox(
                    slide, left=left_m, top=cur_top,
                    width=usable_w, height=Inches(0.7),
                    text=s["subtitle"], font=body_font,
                    size_pt=18, color=muted,
                )
                cur_top = Emu(cur_top + Inches(0.7))
            if s.get("heading"):
                _add_textbox(
                    slide, left=left_m, top=cur_top,
                    width=usable_w, height=Inches(0.8),
                    text=s["heading"], font=display_font,
                    size_pt=22, color=ink, bold=True,
                )
                cur_top = Emu(cur_top + Inches(0.9))
            if s.get("body_md"):
                _add_textbox(
                    slide, left=left_m, top=cur_top,
                    width=usable_w, height=Inches(3.0),
                    text=s["body_md"], font=body_font,
                    size_pt=18, color=ink,
                )
                cur_top = Emu(cur_top + Inches(1.5))
            if s.get("bullets"):
                _add_bullets(
                    slide, left=left_m, top=cur_top,
                    width=usable_w, height=Inches(3.0),
                    bullets=s["bullets"],
                    font=body_font, size_pt=18, color=ink,
                )

        # Footer + slide number on non-cover, non-close slides
        if footer_enabled and stype not in ("cover", "close"):
            _add_textbox(
                slide,
                left=Inches(0.8), top=Inches(7.0),
                width=Inches(6.0), height=Inches(0.4),
                text=company_name,
                font=body_font, size_pt=10, color=muted,
            )
            if slide_numbers:
                _add_textbox(
                    slide,
                    left=Inches(11.5), top=Inches(7.0),
                    width=Inches(1.2), height=Inches(0.4),
                    text=f"{idx} / {len(slides)}",
                    font=body_font, size_pt=10, color=muted,
                    align=PP_ALIGN.RIGHT,
                )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def render_deck_pptx(artifact_id: int) -> dict[str, Any]:
    """Render a deck artifact to .pptx and upload to R2.

    Returns: dict with `r2_key`, `signed_url`, `byte_size`.
    Updates the artifact row's `r2_pptx_key`.
    """
    artifact = await get_artifact(artifact_id)
    if not artifact:
        raise FileNotFoundError(f"artifact #{artifact_id} not found")
    if artifact["kind"] != "deck":
        raise ValueError(
            f"artifact #{artifact_id} is kind={artifact['kind']!r}, not 'deck'"
        )

    kit = load_kit(artifact["business_slug"])
    pptx_bytes = _build_pptx(artifact["content"] or {}, kit)

    key = f"creators/{artifact['business_slug']}/decks/{artifact_id:06d}.pptx"
    s3 = _r2_client()
    s3.put_object(
        Bucket=_r2_bucket(),
        Key=key,
        Body=pptx_bytes,
        ContentType=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
    )
    await update_artifact_render_key(artifact_id, kind="pptx", key=key)

    signed_url = s3.generate_presigned_url(
        "get_object",
        Params={"Bucket": _r2_bucket(), "Key": key},
        ExpiresIn=7 * 24 * 3600,
    )

    return {
        "artifact_id": artifact_id,
        "r2_key": key,
        "signed_url": signed_url,
        "byte_size": len(pptx_bytes),
    }
